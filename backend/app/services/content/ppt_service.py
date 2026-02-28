import os
import uuid
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from app.services.ai.content_planner import plan_presentation


# Output directory for generated files
OUTPUT_DIR = "generated_files"
os.makedirs(OUTPUT_DIR, exist_ok=True)


def hex_to_rgb(hex_color: str):
    """Convert hex color string to RGBColor"""
    hex_color = hex_color.lstrip('#')
    r, g, b = tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
    return RGBColor(r, g, b)


def create_title_slide(prs: Presentation, title: str, subtitle: str, style: str):
    """Create the first title slide with styling"""
    slide_layout = prs.slide_layouts[0]  # Title Slide layout
    slide = prs.slides.add_slide(slide_layout)

    # Color themes per style
    themes = {
        "academic": {"bg": "1E3A5F", "title": "FFFFFF", "sub": "B8D4F0"},
        "simple": {"bg": "FFFFFF", "title": "2C3E50", "sub": "7F8C8D"},
        "technical": {"bg": "0D1117", "title": "58A6FF", "sub": "8B949E"},
        "creative": {"bg": "6C3483", "title": "FFFFFF", "sub": "D7BDE2"},
    }

    theme = themes.get(style, themes["academic"])

    # Set background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(theme["bg"])

    # Title
    title_shape = slide.shapes.title
    title_shape.text = title
    title_tf = title_shape.text_frame.paragraphs[0]
    title_tf.font.size = Pt(40)
    title_tf.font.bold = True
    title_tf.font.color.rgb = hex_to_rgb(theme["title"])
    title_tf.alignment = PP_ALIGN.CENTER

    # Subtitle
    if slide.placeholders[1]:
        sub = slide.placeholders[1]
        sub.text = subtitle
        sub_tf = sub.text_frame.paragraphs[0]
        sub_tf.font.size = Pt(20)
        sub_tf.font.color.rgb = hex_to_rgb(theme["sub"])
        sub_tf.alignment = PP_ALIGN.CENTER

    return slide


def create_content_slide(
    prs: Presentation,
    title: str,
    bullet_points: list,
    speaker_notes: str,
    style: str,
    slide_num: int
):
    """Create a content slide with title and bullet points"""
    slide_layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)

    themes = {
        "academic": {"bg": "F5F8FF", "title": "1E3A5F", "bullet": "2C3E50", "accent": "2E75B6"},
        "simple": {"bg": "FFFFFF", "title": "2C3E50", "bullet": "34495E", "accent": "3498DB"},
        "technical": {"bg": "161B22", "title": "58A6FF", "bullet": "C9D1D9", "accent": "58A6FF"},
        "creative": {"bg": "FDF8FF", "title": "6C3483", "bullet": "4A235A", "accent": "A569BD"},
    }

    theme = themes.get(style, themes["academic"])

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(theme["bg"])

    # Title
    title_shape = slide.shapes.title
    title_shape.text = title
    title_tf = title_shape.text_frame.paragraphs[0]
    title_tf.font.size = Pt(28)
    title_tf.font.bold = True
    title_tf.font.color.rgb = hex_to_rgb(theme["title"])

    # Bullet points
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    for i, point in enumerate(bullet_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = f"• {point}"
        p.font.size = Pt(18)
        p.font.color.rgb = hex_to_rgb(theme["bullet"])
        p.space_before = Pt(6)
        p.level = 0

    # Slide number (bottom right)
    from pptx.util import Inches
    txBox = slide.shapes.add_textbox(Inches(8.5), Inches(6.8), Inches(1), Inches(0.3))
    tf_num = txBox.text_frame
    tf_num.text = str(slide_num)
    tf_num.paragraphs[0].font.size = Pt(12)
    tf_num.paragraphs[0].font.color.rgb = hex_to_rgb(theme["accent"])

    # Speaker notes
    if speaker_notes:
        notes_slide = slide.notes_slide
        notes_tf = notes_slide.notes_text_frame
        notes_tf.text = speaker_notes

    return slide


def generate_ppt(
    topic: str,
    num_slides: int = 8,
    style: str = "academic",
    key_points: list = None,
    subject: str = None,
    include_speaker_notes: bool = True,
    student_name: str = None
) -> dict:
    """
    Complete PPT generation pipeline.

    1. LangChain plans slide structure
    2. python-pptx builds the file
    3. Saves to disk
    4. Returns file path + preview
    """

    print(f"🎨 Generating PPT: {topic}")

    # Step 1 — Plan slides with LangChain + Gemini
    slides_data = plan_presentation(
        topic=topic,
        num_slides=num_slides,
        style=style,
        key_points=key_points,
        subject=subject
    )

    # Step 2 — Create PowerPoint
    prs = Presentation()

    # Set slide dimensions (16:9 widescreen)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Step 3 — Add title slide
    create_title_slide(
        prs=prs,
        title=topic,
        subtitle=f"{subject or 'Presentation'} | {student_name or 'EduAI Suite'}",
        style=style
    )

    # Step 4 — Add content slides
    for i, slide_data in enumerate(slides_data):
        if i == 0:
            continue  # skip first — already used as title

        create_content_slide(
            prs=prs,
            title=slide_data.get("title", f"Slide {i+1}"),
            bullet_points=slide_data.get("bullet_points", []),
            speaker_notes=slide_data.get("speaker_notes", "") if include_speaker_notes else "",
            style=style,
            slide_num=i + 1
        )

    # Step 5 — Save file
    file_id = str(uuid.uuid4())
    file_name = f"{file_id}_{topic[:30].replace(' ', '_')}.pptx"
    file_path = os.path.join(OUTPUT_DIR, file_name)
    prs.save(file_path)

    print(f"✅ PPT saved: {file_path}")

    return {
        "file_path": file_path,
        "file_name": file_name,
        "slides_preview": slides_data,
        "num_slides": len(slides_data)
    }