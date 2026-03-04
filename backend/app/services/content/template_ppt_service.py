import os
import uuid
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import requests
from io import BytesIO

from app.schemas.template_schema import CollegeTemplate
from app.services.ai.content_planner import plan_presentation

OUTPUT_DIR = "generated_files"
os.makedirs(OUTPUT_DIR, exist_ok=True)


def hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert hex string to RGBColor"""
    hex_color = hex_color.lstrip('#')
    r, g, b = tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
    return RGBColor(r, g, b)


def _darken_hex(hex_color: str, factor: float = 0.7) -> str:
    """
    Darken a hex color for text contrast.
    factor < 1 darkens, factor > 1 lightens.
    """
    hex_color = hex_color.lstrip('#')
    r = int(int(hex_color[0:2], 16) * factor)
    g = int(int(hex_color[2:4], 16) * factor)
    b = int(int(hex_color[4:6], 16) * factor)
    return f"{r:02X}{g:02X}{b:02X}"


def _add_logo_to_slide(slide, logo_url: str, position: str = "top_right"):
    """
    Download and add college logo to a slide.
    position: top_right / top_left / bottom_right
    """
    try:
        response = requests.get(logo_url, timeout=5)
        if response.status_code == 200:
            logo_bytes = BytesIO(response.content)

            positions = {
                "top_right": (Inches(8.8), Inches(0.1), Inches(1.0), Inches(0.6)),
                "top_left":  (Inches(0.1), Inches(0.1), Inches(1.0), Inches(0.6)),
                "bottom_right": (Inches(8.8), Inches(6.9), Inches(1.0), Inches(0.5)),
            }

            left, top, width, height = positions.get(
                position, positions["top_right"]
            )
            slide.shapes.add_picture(
                logo_bytes, left, top, width, height
            )
            print("✅ Logo added to slide")
    except Exception as e:
        print(f"Logo add failed: {e}")


def _add_footer_to_slide(slide, template: CollegeTemplate, slide_num: int, total: int):
    """Add college footer bar to bottom of slide"""
    from pptx.util import Emu

    # Footer background bar
    footer_box = slide.shapes.add_textbox(
        Inches(0), Inches(7.1),
        Inches(10), Inches(0.4)
    )
    footer_tf = footer_box.text_frame
    footer_tf.word_wrap = False

    footer_p = footer_tf.paragraphs[0]
    footer_p.alignment = PP_ALIGN.CENTER

    footer_text = f"{template.college_name}"
    if template.footer_text:
        footer_text += f"  |  {template.footer_text}"
    footer_text += f"  |  {slide_num}/{total}"

    run = footer_p.add_run()
    run.text = footer_text
    run.font.size = Pt(9)
    run.font.color.rgb = hex_to_rgb(template.secondary_color)
    run.font.name = template.font_name


def _add_watermark(slide, text: str, template: CollegeTemplate):
    """Add subtle diagonal watermark text"""
    try:
        txBox = slide.shapes.add_textbox(
            Inches(2), Inches(2.5),
            Inches(6), Inches(2)
        )
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER

        run = p.add_run()
        run.text = text[:20]
        run.font.size = Pt(40)
        run.font.bold = True

        # Very light color for watermark
        r = int(int(template.primary_color.lstrip('#')[0:2], 16))
        g = int(int(template.primary_color.lstrip('#')[2:4], 16))
        b = int(int(template.primary_color.lstrip('#')[4:6], 16))
        run.font.color.rgb = RGBColor(
            min(255, r + 200),
            min(255, g + 200),
            min(255, b + 200)
        )

        # Rotate text box (approximate diagonal)
        txBox.rotation = 315
    except Exception as e:
        print(f"Watermark error: {e}")


def create_college_title_slide(
    prs: Presentation,
    topic: str,
    subject: str,
    student_name: str,
    template: CollegeTemplate
):
    """Create title slide with full college branding"""
    slide_layout = prs.slide_layouts[6]  # Blank layout — full control
    slide = prs.slides.add_slide(slide_layout)

    # ── Background ──────────────────
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(template.primary_color)

    # ── Top header banner ──────────────────
    header_box = slide.shapes.add_textbox(
        Inches(0.3), Inches(0.2),
        Inches(9.4), Inches(0.8)
    )
    hf = header_box.text_frame
    hp = hf.paragraphs[0]
    hp.alignment = PP_ALIGN.CENTER
    hr = hp.add_run()
    hr.text = template.header_text or template.college_name
    hr.font.size = Pt(22)
    hr.font.bold = True
    hr.font.name = template.font_name
    hr.font.color.rgb = RGBColor(255, 255, 255)

    # ── Divider line ──────────────────
    line_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.1),
        Inches(9), Inches(0.05)
    )
    lf = line_box.text_frame
    lp = lf.paragraphs[0]
    lr = lp.add_run()
    lr.text = "─" * 80
    lr.font.size = Pt(6)
    lr.font.color.rgb = hex_to_rgb(template.secondary_color)

    # ── Main Title ──────────────────
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.8),
        Inches(9), Inches(2.0)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = topic
    r.font.size = Pt(38)
    r.font.bold = True
    r.font.name = template.font_name
    r.font.color.rgb = RGBColor(255, 255, 255)

    # ── Subject label ──────────────────
    if subject:
        sub_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(3.9),
            Inches(9), Inches(0.5)
        )
        sf = sub_box.text_frame
        sp = sf.paragraphs[0]
        sp.alignment = PP_ALIGN.CENTER
        sr = sp.add_run()
        sr.text = f"Subject: {subject}"
        sr.font.size = Pt(18)
        sr.font.name = template.font_name
        sr.font.color.rgb = hex_to_rgb(template.secondary_color)

    # ── Student name ──────────────────
    if student_name:
        name_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(4.6),
            Inches(9), Inches(0.5)
        )
        nf = name_box.text_frame
        np_ = nf.paragraphs[0]
        np_.alignment = PP_ALIGN.CENTER
        nr = np_.add_run()
        nr.text = f"Presented by: {student_name}"
        nr.font.size = Pt(15)
        nr.font.name = template.font_name
        nr.font.color.rgb = RGBColor(200, 220, 255)

    # ── Footer ──────────────────
    _add_footer_to_slide(slide, template, 1, 1)

    # ── Logo ──────────────────
    if template.logo_url:
        _add_logo_to_slide(slide, template.logo_url, "top_right")


def create_college_content_slide(
    prs: Presentation,
    title: str,
    bullet_points: list,
    speaker_notes: str,
    template: CollegeTemplate,
    slide_num: int,
    total_slides: int,
    add_watermark: bool = True
):
    """Create content slide with college branding"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # ── Background — light version ──────────────────
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(248, 250, 255)

    # ── Top color bar ──────────────────
    top_bar = slide.shapes.add_textbox(
        Inches(0), Inches(0),
        Inches(10), Inches(0.7)
    )
    tb_fill = top_bar.fill
    tb_fill.solid()
    tb_fill.fore_color.rgb = hex_to_rgb(template.primary_color)

    # College name in top bar
    tbf = top_bar.text_frame
    tbp = tbf.paragraphs[0]
    tbp.alignment = PP_ALIGN.LEFT
    tbr = tbp.add_run()
    tbr.text = f"  {template.college_name}"
    tbr.font.size = Pt(11)
    tbr.font.name = template.font_name
    tbr.font.color.rgb = RGBColor(255, 255, 255)

    # ── Slide title ──────────────────
    title_box = slide.shapes.add_textbox(
        Inches(0.4), Inches(0.8),
        Inches(9.2), Inches(0.9)
    )
    ttf = title_box.text_frame
    ttp = ttf.paragraphs[0]
    ttr = ttp.add_run()
    ttr.text = title
    ttr.font.size = Pt(28)
    ttr.font.bold = True
    ttr.font.name = template.font_name
    ttr.font.color.rgb = hex_to_rgb(template.primary_color)

    # ── Title underline ──────────────────
    ul_box = slide.shapes.add_textbox(
        Inches(0.4), Inches(1.65),
        Inches(9.2), Inches(0.05)
    )
    ulf = ul_box.text_frame
    ulp = ulf.paragraphs[0]
    ulr = ulp.add_run()
    ulr.text = "─" * 100
    ulr.font.size = Pt(5)
    ulr.font.color.rgb = hex_to_rgb(template.secondary_color)

    # ── Bullet points ──────────────────
    content_box = slide.shapes.add_textbox(
        Inches(0.4), Inches(1.9),
        Inches(9.2), Inches(4.8)
    )
    ctf = content_box.text_frame
    ctf.word_wrap = True

    for i, point in enumerate(bullet_points):
        if i == 0:
            cp = ctf.paragraphs[0]
        else:
            cp = ctf.add_paragraph()

        cp.space_before = Pt(8)
        cr = cp.add_run()
        cr.text = f"  ▸  {point}"
        cr.font.size = Pt(17)
        cr.font.name = template.font_name
        cr.font.color.rgb = RGBColor(40, 40, 60)

    # ── Watermark ──────────────────
    if add_watermark and template.watermark_text:
        _add_watermark(slide, template.watermark_text, template)

    # ── Footer ──────────────────
    _add_footer_to_slide(slide, template, slide_num, total_slides)

    # ── Logo on content slides ──────────────────
    if template.logo_url:
        _add_logo_to_slide(slide, template.logo_url, "top_right")

    # ── Speaker notes ──────────────────
    if speaker_notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = speaker_notes


def generate_college_ppt(
    topic: str,
    template: CollegeTemplate,
    num_slides: int = 8,
    key_points: list = None,
    subject: str = None,
    include_speaker_notes: bool = True,
    student_name: str = None
) -> dict:
    """
    Generate a PPT using college template branding.
    Main entry point called from the route.
    """
    print(f"🎨 Generating college-branded PPT: {topic}")
    print(f"   Template: {template.college_name}")
    print(f"   Colors: {template.primary_color} / {template.secondary_color}")

    # Step 1 — Plan slides with AI
    slides_data = plan_presentation(
        topic=topic,
        num_slides=num_slides,
        style="academic",
        key_points=key_points,
        subject=subject
    )

    total = len(slides_data)

    # Step 2 — Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Step 3 — Title slide
    create_college_title_slide(
        prs=prs,
        topic=topic,
        subject=subject,
        student_name=student_name,
        template=template
    )

    # Step 4 — Content slides
    for i, slide_data in enumerate(slides_data):
        if i == 0:
            continue  # skip first — already made as title

        create_college_content_slide(
            prs=prs,
            title=slide_data.get("title", f"Slide {i+1}"),
            bullet_points=slide_data.get("bullet_points", []),
            speaker_notes=slide_data.get("speaker_notes", "") if include_speaker_notes else "",
            template=template,
            slide_num=i + 1,
            total_slides=total,
            add_watermark=(i == 1)  # watermark only on slide 2
        )

    # Step 5 — Save file
    file_id = str(uuid.uuid4())
    file_name = f"{file_id}_{topic[:25].replace(' ', '_')}_CollegeTemplate.pptx"
    file_path = os.path.join(OUTPUT_DIR, file_name)
    prs.save(file_path)

    print(f"✅ College PPT saved: {file_name}")

    return {
        "file_path": file_path,
        "file_name": file_name,
        "slides_preview": slides_data,
        "num_slides": total,
        "template_used": template.college_name
    }